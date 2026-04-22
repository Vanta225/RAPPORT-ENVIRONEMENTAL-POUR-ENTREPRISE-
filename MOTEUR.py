from pptx import Presentation
from pptx.util import Inches
import os

def generer_rapport_audit(donnees_formulaire, photos_dict, template_path, output_path):
    """
    donnees_formulaire: dict contenant les textes { 'INS_CODE': 'ABJ001', ... }
    photos_dict: dict contenant les objets images { 'INS_SITE1': file_object, ... }
    """
    prs = Presentation(template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            # --- TRAITEMENT DES TEXTES ---
            if shape.has_text_frame:
                for balise, valeur in donnees_formulaire.items():
                    if shape.text == balise:
                        shape.text = str(valeur) if valeur else ""

            # --- TRAITEMENT DES IMAGES ---
            # On vérifie si le nom de la forme est dans notre dictionnaire de photos
            if shape.name in photos_dict and photos_dict[shape.name] is not None:
                photo = photos_dict[shape.name]
                
                # Récupération des dimensions de la zone définie dans PowerPoint
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # Insertion de la photo à l'emplacement exact
                slide.shapes.add_picture(photo, left, top, width, height)

                # Suppression du rectangle de repère (optionnel mais plus propre)
                sp = shape._element
                sp.getparent().remove(sp)

    prs.save(output_path)
    return output_path
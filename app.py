import streamlit as st
from pptx import Presentation
import io
from datetime import datetime
from PIL import Image

# --- CONFIGURATION ---
MOT_DE_PASSE_ENTREPRISE = "ETECHNOLOGIE"

# --- FONCTION UPLOADER AVEC PREVIEW + EDITION ---
def uploader_avec_preview_avance(label, key, types=['jpg','png','jpeg']):
    fichier = st.file_uploader(label, type=types, key=key)

    if fichier is not None:
        image = Image.open(fichier)

        st.markdown("🔧 Ajustement")

        col1, col2 = st.columns(2)

        rotation = col1.selectbox(
            "Rotation",
            [0, 90, 180, 270],
            key=f"rot_{key}"
        )

        zoom = col2.slider(
            "Zoom aperçu",
            50, 300, 150,
            key=f"zoom_{key}"
        )

        image_modifiee = image.rotate(rotation, expand=True)

        st.markdown("👁️ Aperçu du rendu")
        st.image(image_modifiee, width=zoom)

        return image_modifiee

    return None

# --- MOTEUR PPT ---
def generer_rapport(donnees, photos, template_path):
    try:
        prs = Presentation(template_path)
    except Exception as e:
        st.error(f"Erreur chargement template : {e}")
        return None
    
    for slide in prs.slides:
        for shape in slide.shapes:

            # TEXTE
            if shape.has_text_frame:
                for balise, valeur in donnees.items():
                    if balise in shape.text:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if balise in run.text:
                                    run.text = run.text.replace(balise, str(valeur) if valeur else "")

            # IMAGES
            if shape.name in photos and photos[shape.name] is not None:
                photo = photos[shape.name]
                left, top, width, height = shape.left, shape.top, shape.width, shape.height

                if isinstance(photo, Image.Image):
                    img_bytes = io.BytesIO()
                    photo.save(img_bytes, format='PNG')
                    img_bytes.seek(0)
                    slide.shapes.add_picture(img_bytes, left, top, width, height)
                else:
                    slide.shapes.add_picture(photo, left, top, width, height)

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

# --- AUTH ---
def est_authentifie():
    if "auth" not in st.session_state:
        st.session_state["auth"] = False

    if not st.session_state["auth"]:
        st.markdown("### 🔒 Accès Sécurisé")
        mdp = st.text_input("Mot de passe", type="password")
        if st.button("Connexion"):
            if mdp == MOT_DE_PASSE_ENTREPRISE:
                st.session_state["auth"] = True
                st.rerun()
            else:
                st.error("Incorrect")
        return False
    return True

# --- APP ---
if est_authentifie():
    st.set_page_config(layout="wide")
    st.title("📱 Audit Site Telecom")

    donnees = {}
    photos = {}
    date_du_jour = datetime.now().strftime("%d-%m-%Y")

    tabs = st.tabs([
        "INFOS", "VUES", "SECURITE",
        "PROPRETE", "PYLONE", "ELEC",
        "GE", "CLIM", "ANOMALIES"
    ])

    # INFOS
    with tabs[0]:
        donnees["INS_NOM"] = st.text_input("Nom site")
        donnees["INS_CODE"] = st.text_input("Code site")

    # VUES
    with tabs[1]:
        photos["INS_VUE_DU_SITE"] = uploader_avec_preview_avance("Vue site", "vue")
        photos["INS_PLAQUE"] = uploader_avec_preview_avance("Plaque", "plaque")

    # SECURITE
    with tabs[2]:
        photos["INS_ACCES"] = uploader_avec_preview_avance("Accès", "acces")

    # PROPRETE
    with tabs[3]:
        photos["INS_PROPRE"] = uploader_avec_preview_avance("Propreté", "propre")

    # PYLONE
    with tabs[4]:
        photos["INS_FONDA1"] = uploader_avec_preview_avance("Fondation", "fonda")

    # ELEC
    with tabs[5]:
        photos["INS_TGBT"] = uploader_avec_preview_avance("TGBT", "tgbt")

    # GE
    with tabs[6]:
        photos["INS_MOTEUR"] = uploader_avec_preview_avance("Moteur GE", "ge")

    # CLIM
    with tabs[7]:
        photos["INS_CLIM1"] = uploader_avec_preview_avance("Clim", "clim")

    # ANOMALIES
    with tabs[8]:
        photos["INS_AN1"] = uploader_avec_preview_avance("Anomalie", "ano")
        donnees["INS_REMARQUES"] = st.text_area("Remarques")

        if st.button("🚀 GENERER RAPPORT"):
            if not donnees["INS_CODE"] or not donnees["INS_NOM"]:
                st.error("Nom et Code obligatoires")
            else:
                output = generer_rapport(donnees, photos, "template.pptx")

                if output:
                    nom = f"RAPPORT_{donnees['INS_CODE']}_{date_du_jour}.pptx"
                    st.download_button("📥 Télécharger", data=output, file_name=nom)
